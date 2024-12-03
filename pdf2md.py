import fitz
import os
import argparse
import logging
from datetime import datetime
from pathlib import Path
from tqdm import tqdm
from docx import Document
from pptx import Presentation
from striprtf.striprtf import rtf_to_text
import mimetypes
import pdfplumber
import pytesseract
from pdf2image import convert_from_path
import tempfile
import shutil

class ConversionReport:
    """Classe para gerenciar o relatório de conversão"""
    def __init__(self):
        self.successful = []
        self.failed = []
        self.skipped = []
        
    def add_success(self, file_path):
        self.successful.append(file_path)
        
    def add_failure(self, file_path, error):
        self.failed.append((file_path, str(error)))
        
    def add_skipped(self, file_path, reason):
        self.skipped.append((file_path, reason))
        
    def generate_report(self, output_dir):
        report_path = Path(output_dir) / f"conversion_report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.txt"
        with open(report_path, 'w', encoding='utf-8') as f:
            f.write("=== Relatório de Conversão ===\n\n")
            
            f.write(f"Total de arquivos processados: {len(self.successful) + len(self.failed) + len(self.skipped)}\n")
            f.write(f"Convertidos com sucesso: {len(self.successful)}\n")
            f.write(f"Falhas na conversão: {len(self.failed)}\n")
            f.write(f"Arquivos ignorados: {len(self.skipped)}\n\n")
            
            if self.successful:
                f.write("\n=== Arquivos Convertidos com Sucesso ===\n")
                for file in self.successful:
                    f.write(f"✓ {file}\n")
            
            if self.failed:
                f.write("\n=== Falhas na Conversão ===\n")
                for file, error in self.failed:
                    f.write(f"✗ {file}\nErro: {error}\n\n")
            
            if self.skipped:
                f.write("\n=== Arquivos Ignorados ===\n")
                for file, reason in self.skipped:
                    f.write(f"- {file}\nMotivo: {reason}\n\n")
        
        return report_path

class DocumentConverter:
    """Classe para converter diferentes tipos de documentos para Markdown"""
    
    def __init__(self, report):
        self.report = report
        # Configura logging
        logging.basicConfig(
            level=logging.INFO,
            format='%(asctime)s - %(levelname)s - %(message)s',
            handlers=[
                logging.StreamHandler(),
                logging.FileHandler('conversion.log')
            ]
        )
        self.logger = logging.getLogger(__name__)
        
        # Registra os tipos MIME
        mimetypes.init()
        mimetypes.add_type('application/rtf', '.rtf')
        mimetypes.add_type('application/vnd.openxmlformats-officedocument.presentationml.presentation', '.pptx')
        
        # Verifica se o Tesseract está instalado
        try:
            pytesseract.get_tesseract_version()
            self.ocr_available = True
        except:
            self.logger.warning("Tesseract não está instalado. OCR não estará disponível.")
            self.ocr_available = False

    def extract_text_with_ocr(self, pdf_path):
        """Extrai texto do PDF usando OCR"""
        self.logger.info("Tentando extrair texto com OCR")
        text_parts = []
        
        try:
            # Converte PDF para imagens
            with tempfile.TemporaryDirectory() as temp_dir:
                images = convert_from_path(pdf_path)
                
                # Processa cada página com OCR
                for i, image in enumerate(images):
                    self.logger.debug(f"Processando página {i+1} com OCR")
                    text = pytesseract.image_to_string(image, lang='por')
                    if text.strip():
                        text_parts.append(text)
                    
            return "\n\n".join(text_parts)
        except Exception as e:
            self.logger.error(f"Erro no OCR: {str(e)}")
            return ""

    def extract_text_with_pdfplumber(self, pdf_path):
        """Extrai texto do PDF usando pdfplumber"""
        self.logger.info("Tentando extrair texto com pdfplumber")
        text_parts = []
        
        try:
            with pdfplumber.open(pdf_path) as pdf:
                for page in pdf.pages:
                    text = page.extract_text()
                    if text and text.strip():
                        text_parts.append(text)
            
            return "\n\n".join(text_parts)
        except Exception as e:
            self.logger.error(f"Erro com pdfplumber: {str(e)}")
            return ""

    def convert_pdf(self, file_path):
        """Converte PDF para texto markdown usando múltiplos métodos"""
        doc = None
        text_content = ""
        
        try:
            # 1. Primeira tentativa: PyMuPDF (mais rápido e preciso)
            try:
                doc = fitz.open(file_path)
                if doc.is_encrypted:
                    if not doc.authenticate(""):
                        raise ValueError("PDF está protegido com senha")
                
                if doc.page_count == 0:
                    raise ValueError("PDF está vazio")
                
                markdown_text = []
                for page_num, page in enumerate(doc, 1):
                    self.logger.debug(f"Processando página {page_num} com PyMuPDF")
                    blocks = page.get_text("dict")["blocks"]
                    for block in blocks:
                        if "lines" in block:
                            for line in block["lines"]:
                                for span in line["spans"]:
                                    text = span["text"].strip()
                                    if text:
                                        if span["size"] > 12:
                                            text = f"## {text}"
                                        markdown_text.append(text)
                    markdown_text.append("\n")
                
                text_content = "\n".join(markdown_text)
                if text_content.strip():
                    return text_content
            except Exception as e:
                self.logger.warning(f"PyMuPDF falhou: {str(e)}")
            
            # 2. Segunda tentativa: pdfplumber
            text_content = self.extract_text_with_pdfplumber(file_path)
            if text_content.strip():
                return text_content
            
            # 3. Última tentativa: OCR
            if self.ocr_available:
                text_content = self.extract_text_with_ocr(file_path)
                if text_content.strip():
                    return text_content
            
            raise ValueError("Não foi possível extrair texto do PDF usando nenhum dos métodos disponíveis")
            
        except Exception as e:
            if "password" in str(e).lower():
                raise ValueError("PDF está protegido com senha")
            elif "xref" in str(e).lower():
                raise ValueError("PDF está corrompido ou tem estrutura inválida")
            else:
                raise Exception(str(e))
        finally:
            if doc:
                doc.close()

    def convert_docx(self, file_path):
        """Converte DOCX para texto markdown"""
        try:
            doc = Document(file_path)
            markdown_text = []
            
            for paragraph in doc.paragraphs:
                if paragraph.style.name.startswith('Heading'):
                    level = paragraph.style.name[-1]
                    markdown_text.append(f"{'#' * int(level)} {paragraph.text}")
                else:
                    markdown_text.append(paragraph.text)
            
            return "\n\n".join(markdown_text)
        except Exception as e:
            raise Exception(f"Erro ao processar DOCX: {str(e)}")

    def convert_pptx(self, file_path):
        """Converte PPTX para texto markdown"""
        try:
            prs = Presentation(file_path)
            markdown_text = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                markdown_text.append(f"## Slide {slide_num}")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        markdown_text.append(shape.text)
                markdown_text.append("\n---\n")
            
            return "\n".join(markdown_text)
        except Exception as e:
            raise Exception(f"Erro ao processar PPTX: {str(e)}")

    def convert_txt(self, file_path):
        """Converte TXT para texto markdown"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                return file.read()
        except Exception as e:
            raise Exception(f"Erro ao processar TXT: {str(e)}")

    def convert_rtf(self, file_path):
        """Converte RTF para texto markdown"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                return rtf_to_text(file.read())
        except Exception as e:
            raise Exception(f"Erro ao processar RTF: {str(e)}")

    def convert_file(self, input_path, output_path):
        """Converte qualquer arquivo suportado para markdown"""
        self.logger.info(f"Iniciando conversão de: {input_path}")
        
        try:
            # Verifica se o arquivo existe
            if not os.path.exists(input_path):
                raise FileNotFoundError(f"Arquivo não encontrado: {input_path}")
            
            # Verifica se o arquivo está vazio
            if os.path.getsize(input_path) == 0:
                raise ValueError("Arquivo está vazio")
            
            # Detecta o tipo do arquivo
            mime_type, _ = mimetypes.guess_type(input_path)
            
            # Converte o arquivo
            with open(output_path, 'w', encoding='utf-8') as out_file:
                if mime_type == 'application/pdf':
                    content = self.convert_pdf(input_path)
                elif mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
                    content = self.convert_docx(input_path)
                elif mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
                    content = self.convert_pptx(input_path)
                elif mime_type == 'application/rtf':
                    content = self.convert_rtf(input_path)
                elif mime_type == 'text/plain' or input_path.endswith('.txt'):
                    content = self.convert_txt(input_path)
                else:
                    self.report.add_skipped(input_path, f"Formato não suportado: {mime_type}")
                    return
                
                # Verifica se o conteúdo convertido não está vazio
                if not content.strip():
                    raise ValueError("Nenhum conteúdo extraído do arquivo")
                
                out_file.write(content)
                self.report.add_success(input_path)
                self.logger.info(f"Arquivo convertido com sucesso: {input_path}")
                
        except Exception as e:
            self.logger.error(f"Erro ao processar {input_path}: {str(e)}")
            self.report.add_failure(input_path, str(e))

def process_directory(input_dir, output_dir):
    """
    Processa recursivamente um diretório convertendo todos os documentos suportados para Markdown.
    """
    input_path = Path(input_dir)
    output_path = Path(output_dir)
    report = ConversionReport()
    converter = DocumentConverter(report)
    
    # Verifica se o diretório de entrada existe
    if not input_path.exists():
        raise FileNotFoundError(f"Diretório de entrada não encontrado: {input_dir}")
    
    # Cria o diretório de saída se não existir
    output_path.mkdir(parents=True, exist_ok=True)
    
    # Lista todos os arquivos no diretório e subdiretórios
    supported_extensions = {'.pdf', '.docx', '.pptx', '.txt', '.rtf'}
    files = [f for f in input_path.rglob("*") if f.suffix.lower() in supported_extensions]
    
    if not files:
        print(f"Nenhum arquivo suportado encontrado em: {input_dir}")
        return
    
    print(f"Encontrados {len(files)} arquivos para processar.")
    
    # Processa cada arquivo com barra de progresso
    for file in tqdm(files, desc="Convertendo documentos"):
        # Mantém a estrutura de diretórios relativa
        relative_path = file.relative_to(input_path)
        output_file = output_path / relative_path.with_suffix('.md')
        
        # Cria os subdiretórios necessários
        output_file.parent.mkdir(parents=True, exist_ok=True)
        
        # Converte o arquivo para Markdown
        converter.convert_file(str(file), str(output_file))
    
    # Gera o relatório
    report_path = report.generate_report(output_dir)
    print(f"\nRelatório de conversão gerado em: {report_path}")
    
    # Exibe um resumo
    print(f"\nResumo da conversão:")
    print(f"✓ Convertidos com sucesso: {len(report.successful)}")
    print(f"✗ Falhas na conversão: {len(report.failed)}")
    print(f"- Arquivos ignorados: {len(report.skipped)}")

def main():
    parser = argparse.ArgumentParser(description='Converte documentos para Markdown otimizado para LLMs')
    parser.add_argument('--input_dir', required=True, help='Diretório de entrada contendo os documentos')
    parser.add_argument('--output_dir', required=True, help='Diretório de saída para os arquivos Markdown')
    
    args = parser.parse_args()
    
    try:
        print(f"Iniciando conversão de documentos em: {args.input_dir}")
        process_directory(args.input_dir, args.output_dir)
        print("Conversão concluída!")
    except Exception as e:
        print(f"Erro fatal durante a execução: {str(e)}")
        exit(1)

if __name__ == "__main__":
    main()
